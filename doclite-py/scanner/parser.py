"""
文件解析模块 - 从不同格式的文件中提取文本内容

支持的文件格式:
- PDF: 使用 PyMuPDF 提取文本
- DOCX: 使用 python-docx 提取段落文本
- PPTX: 使用 python-pptx 提取幻灯片文本
- XLSX: 使用 openpyxl 提取单元格文本
- MD/TXT: 直接读取纯文本内容
- Image: 使用 Tesseract OCR 提取图片文字（需要 pytesseract 和 Pillow）
"""

import os
import logging
from typing import Dict, Callable, Any
import fitz  # PyMuPDF
from docx import Document

logger = logging.getLogger(__name__)

# 文件类型到提取函数的映射
EXTRACTORS: Dict[str, Callable[[str], str]] = {}

def _register_extractor(file_type: str):
    """注册文件类型提取器的装饰器"""
    def decorator(func: Callable[[str], str]):
        EXTRACTORS[file_type] = func
        return func
    return decorator

@_register_extractor("pdf")
def _extract_pdf(file_path: str) -> str:
    """提取 PDF 文件文本（含嵌入图片 OCR）"""
    text_parts = []
    with fitz.open(file_path) as doc:
        for page in doc:
            text_parts.append(page.get_text())
            image_text = _ocr_pdf_page_images(doc, page)
            if image_text:
                text_parts.append("\n[图片OCR]\n" + image_text + "\n")
    return "".join(text_parts).strip()


def _ocr_pdf_page_images(doc, page) -> str:
    """对 PDF 单页的嵌入图片做 OCR，返回拼接文本；失败/禁用返回空串"""
    try:
        import pytesseract
    except ImportError:
        logger.warning("pytesseract 未安装，跳过 PDF 图片 OCR")
        return ""

    from api.settings import load_settings
    settings = load_settings()
    if not settings.ocr_enabled:
        return ""

    image_texts = []
    for img_info in page.get_images(full=True):
        xref = img_info[0]
        try:
            pix = fitz.Pixmap(doc, xref)
            if pix.n - pix.alpha >= 4:
                pix = fitz.Pixmap(fitz.csRGB, pix)
            import tempfile
            tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            tmp_path = tmp.name
            tmp.close()
            try:
                pix.save(tmp_path)
                text = pytesseract.image_to_string(tmp_path, lang=settings.ocr_language).strip()
                if text:
                    image_texts.append(text)
            finally:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
        except Exception as e:
            logger.warning(f"PDF 图片 OCR 失败 xref={xref}: {e}")
            continue
    return "\n".join(image_texts)

@_register_extractor("docx")
def _extract_docx(file_path: str) -> str:
    """提取 Word 文档文本"""
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs]).strip()

@_register_extractor("pptx")
def _extract_pptx(file_path: str) -> str:
    """提取 PowerPoint 文件文本"""
    from pptx import Presentation
    
    prs = Presentation(file_path)
    text_parts = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
    
    return "\n".join(text_parts).strip()

@_register_extractor("xlsx")
def _extract_xlsx(file_path: str) -> str:
    """提取 Excel 文件文本"""
    from openpyxl import load_workbook
    
    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []
    
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text_parts.append(f"[工作表: {sheet}]")
        
        for row in ws.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else "" for cell in row]
            row_str = "\t".join(row_text)
            if row_str.strip():
                text_parts.append(row_str)
    
    wb.close()
    return "\n".join(text_parts).strip()

@_register_extractor("md")
@_register_extractor("txt")
def _extract_plain_text(file_path: str) -> str:
    """提取纯文本文件内容"""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()

@_register_extractor("png")
@_register_extractor("jpg")
@_register_extractor("jpeg")
@_register_extractor("bmp")
@_register_extractor("tiff")
@_register_extractor("tif")
@_register_extractor("webp")
@_register_extractor("gif")
def _extract_image(file_path: str) -> str:
    """提取图片文件文本（OCR）"""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        logger.warning("pytesseract 或 Pillow 未安装，无法进行 OCR 识别")
        return ""
    
    from api.settings import load_settings
    settings = load_settings()
    
    if not settings.ocr_enabled:
        return ""
    
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang=settings.ocr_language)
        return text.strip()
    except Exception as e:
        logger.error(f"OCR 识别失败 {file_path}: {e}")
        return ""

def extract_text(file_info: Dict[str, Any]) -> str:
    """
    根据文件类型提取纯文本内容
    
    Args:
        file_info: 文件信息字典，包含 path, file_type 等字段
    
    Returns:
        提取的文本内容，失败返回空字符串
    """
    path: str = file_info["path"]
    file_type: str = file_info["file_type"]

    try:
        extractor = EXTRACTORS.get(file_type)
        if extractor:
            return extractor(path)
        else:
            logger.warning(f"不支持的文件类型: {file_type}")
            return ""
    except Exception as e:
        logger.error(f"解析文件失败 {path}: {e}")
        return ""
